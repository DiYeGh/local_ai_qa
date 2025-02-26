import os
import re
import sys
import tempfile
import uuid
import fitz  # PyMuPDF 用于处理PDF文件
import docx  # python-docx 用于处理Word文件
from pptx import Presentation  # python-pptx 用于处理PPT文件
from win32com import client as win32
from docx.document import Document as _Document
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.table import _Cell, Table, _Row
from docx.text.paragraph import Paragraph


#
def extract_text_from_pdf(pdf_path):
    # 定义页眉页脚的高度占比
    header_footer_threshold = 0.1
    # 打开pdf文件
    document = fitz.open(pdf_path)
    extracted_text = ""
    for page_num in range(len(document)):
        page = document.load_page(page_num)
        # 获取页面的文本块
        text_blocks = page.get_text("blocks")
        # 获取页面高度
        page_height = page.rect.height
        page_text = ""
        for block in text_blocks:
            x0, y0, x1, y1, text = block[:5]  # 文本块位置和内容
            # 如果文本位于页眉或页脚区域，跳过处理
            if y0 < header_footer_threshold * page_height or y1 > (1 - header_footer_threshold) * page_height:
                continue
            page_text += text.strip() + "\n"
        extracted_text += page_text.strip() + "\n\n"
    return extracted_text


# 提取ppt、pptx中的文本
def extract_text_from_pptx(ppt_path):
    # 打开PPT文件
    presentation = Presentation(ppt_path)
    text_content = []
    # 遍历每一张幻灯片
    for slide in presentation.slides:
        # 遍历幻灯片中的每一个形状
        for shape in slide.shapes:
            # 判断形状是否有文本框
            if hasattr(shape, "text"):
                # 提取文本框中的文字
                text_content.append(shape.text)
            # 检查形状是否是表格
            if shape.has_table:
                table = shape.table
                # 遍历表格的行和列，提取每个单元格的文本
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text)
                    # 使用 \t 将每行的单元格内容连接成一行
                    text_content.append("\t".join(row_text))
    return "\n".join(text_content)


# 提取Word文件中的文本
def extract_text_from_docx(docx_path):
    document = docx.Document(docx_path)

    def iterate_block_items(parent):
        """递归遍历文档中的段落和表格"""
        if isinstance(parent, _Document):
            parent_element = parent.element.body
        elif isinstance(parent, _Cell):
            parent_element = parent._tc
        elif isinstance(parent, _Row):
            parent_element = parent._tr
        else:
            raise ValueError("文档结构不正确")
        for child in parent_element.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    def is_cell_merged_vertically(cell):
        """判断单元格是否为垂直合并的一部分，且不是合并的起始单元格"""
        tc = cell._tc
        vMerge_val = tc.xpath(".//w:vMerge/@w:val")
        return bool(vMerge_val) and vMerge_val[0] != "restart"

    def is_cell_merged_horizontally(cell):
        """判断单元格是否为水平合并的一部分"""
        tc = cell._tc
        gridSpan_val = tc.xpath(".//w:gridSpan/@w:val")
        return bool(gridSpan_val) and int(gridSpan_val[0]) > 1

    # 缓存机制，记录已经处理的合并单元格
    processed_cells = set()
    extracted_content = []
    for block in iterate_block_items(document):
        # 如果是段落，提取文本
        if isinstance(block, Paragraph):
            paragraph_text = block.text.strip()
            if paragraph_text:
                extracted_content.append(paragraph_text)
        # 如果是表格，处理表格中的每一行
        elif isinstance(block, Table):
            for row in block.rows:
                row_data = []
                for col_idx, cell in enumerate(row.cells):
                    cell_key = (row._index, col_idx)
                    # 如果该单元格已经被处理过（垂直或水平合并），跳过
                    if cell_key in processed_cells:
                        continue
                    # 垂直合并单元格处理
                    if is_cell_merged_vertically(cell):
                        continue
                    # 水平合并单元格处理，记录起始单元格，跳过后续单元格
                    if is_cell_merged_horizontally(cell):
                        # 把该单元格之后的所有被合并的单元格都标记为已处理
                        for i in range(1, int(cell._tc.xpath(".//w:gridSpan/@w:val")[0])):
                            processed_cells.add((row._index, col_idx + i))
                    # 提取单元格文本
                    cell_text = cell.text.strip()
                    row_data.append(cell_text)
                    # 记录已经处理过的单元格
                    processed_cells.add(cell_key)
                # 只添加包含文本的行
                if row_data:
                    extracted_content.append("\t".join(row_data))  # 使用制表符对齐表格内容
    return "\n".join(extracted_content)


def doc_to_docx(doc_path, docx_path):
    try:
        word_app = win32.Dispatch('Word.Application')
        doc = word_app.Documents.Open(doc_path)
        doc.SaveAs(docx_path, FileFormat=16)  # 16 表示 .docx 格式
        doc.Close()
        word_app.Quit()
    except Exception as e:
        print(f"Error while converting document: {e}")
        word_app.Quit()


def ppt_to_pptx(ppt_path, pptx_path):
    try:
        powerpoint_app = win32.Dispatch('PowerPoint.Application')
        ppt = powerpoint_app.Presentations.Open(ppt_path)
        ppt.SaveAs(pptx_path, FileFormat=24)  # 24 表示 .pptx 格式
        ppt.Close()
        powerpoint_app.Quit()
    except Exception as e:
        print(f"Error while converting presentation: {e}")
        powerpoint_app.Quit()


def extract_text_from_txt(txt_path):
    with open(txt_path, 'r', encoding='utf-8') as file:
        text = file.read()
    return text


def extract_text_from_doc(doc_path):
    # 创建一个以UUID命名的临时文件
    temp_docx_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.docx")
    # 将doc文件临时转换存储到临时docx文件
    doc_to_docx(doc_path, temp_docx_path)
    # 获取文件文字
    text = extract_text_from_docx(temp_docx_path)
    # 删除临时文件
    os.remove(temp_docx_path)
    return text


def extract_text_from_ppt(ppt_path):
    # 创建一个以UUID命名的临时文件
    temp_pptx_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.pptx")
    # 将ppt文件临时转换存储到临时pptx文件
    ppt_to_pptx(ppt_path, temp_pptx_path)
    # 获取文件文字
    text = extract_text_from_pptx(temp_pptx_path)
    # 删除临时文件
    os.remove(temp_pptx_path)
    return text


# 通用文件文本提取函数
def extract_text_from_file(file_path):
    not_support = []
    text = ''
    # 获取文件的扩展名并转换为小写
    file_extension = os.path.splitext(file_path)[1].lower()
    if file_extension == '.docx':
        text = extract_text_from_docx(file_path)
    elif file_extension == '.pdf':
        text = extract_text_from_pdf(file_path)
    elif file_extension == '.doc' or file_extension == '.docm':
        text = extract_text_from_doc(file_path)
    elif file_extension == '.ppt':
        text = extract_text_from_ppt(file_path)
    elif file_extension == '.pptx':
        text = extract_text_from_pptx(file_path)
    elif file_extension == '.txt':
        text = extract_text_from_txt(file_path)
    else:
        not_support.append(file_path)
    return text